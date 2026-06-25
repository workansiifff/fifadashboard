import os
import sys

# Define target pptx path
pptx_path = "FIFA_World_Cup_2026_Dashboard_Presentation.pptx"

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx is not installed. Attempting to install...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette (FIFA Navy Theme)
COLOR_BG = RGBColor(2, 6, 23)        # Deep Navy (#020617)
COLOR_CARD = RGBColor(10, 20, 42)     # Glass card background (#0a142a)
COLOR_TEXT_MAIN = RGBColor(248, 250, 252) # White (#f8fafc)
COLOR_TEXT_MUTED = RGBColor(148, 163, 184) # Muted Gray (#94a3b8)
COLOR_ACCENT = RGBColor(56, 189, 248)  # Light Blue (#38bdf8)
COLOR_GOLD = RGBColor(251, 191, 36)    # Gold (#fbbf24)

def apply_background(slide):
    # Add a rectangle that covers the entire slide for the dark background
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLOR_BG
    rect.line.fill.background() # No border
    
    # Add top decorative accent line (FIFA Gold/Blue gradient look)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.1)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_ACCENT
    accent.line.fill.background()

def create_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    
    # Decorative glass-like panel in center
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLOR_CARD
    panel.line.color.rgb = COLOR_ACCENT
    panel.line.width = Pt(1.5)
    
    # Text Frame inside the panel
    txBox = slide.shapes.add_textbox(Inches(2.0), Inches(2.0), Inches(9.333), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Badge
    p0 = tf.paragraphs[0]
    p0.text = "🏆 FIFA WORLD CUP 2026"
    p0.font.name = 'Outfit'
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_GOLD
    p0.alignment = PP_ALIGN.CENTER
    p0.space_after = Pt(24)
    
    # Title
    p1 = tf.add_paragraph()
    p1.text = "Tournament Analytics Dashboard"
    p1.font.name = 'Outfit'
    p1.font.size = Pt(38)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_MAIN
    p1.alignment = PP_ALIGN.CENTER
    p1.space_after = Pt(10)
    
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "Premium Front-End Project Presentation"
    p2.font.name = 'Plus Jakarta Sans'
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_ACCENT
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(40)
    
    # Team Info
    p3 = tf.add_paragraph()
    p3.text = "Team: FIFO  |  Country: Uruguay  |  Developer: Uruguay Developer"
    p3.font.name = 'Plus Jakarta Sans'
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_TEXT_MUTED
    p3.alignment = PP_ALIGN.CENTER

def create_content_slide(title, subtitle, points_left, points_right=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_background(slide)
    
    # Header Area
    header_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.833), Inches(1.2))
    tf_header = header_box.text_frame
    tf_header.word_wrap = True
    
    # Section tag
    p_tag = tf_header.paragraphs[0]
    p_tag.text = subtitle.upper()
    p_tag.font.name = 'Outfit'
    p_tag.font.size = Pt(12)
    p_tag.font.bold = True
    p_tag.font.color.rgb = COLOR_ACCENT
    p_tag.space_after = Pt(4)
    
    # Title
    p_title = tf_header.add_paragraph()
    p_title.text = title
    p_title.font.name = 'Outfit'
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_MAIN
    
    # Two Columns or One Column Layout
    if points_right:
        # Left Box
        left_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
        tf_left = left_box.text_frame
        tf_left.word_wrap = True
        add_bullets(tf_left, points_left)
        
        # Right Box (Decorative panel)
        right_panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.75), Inches(1.8), Inches(5.8), Inches(4.8)
        )
        right_panel.fill.solid()
        right_panel.fill.fore_color.rgb = COLOR_CARD
        right_panel.line.color.rgb = COLOR_TEXT_MUTED
        right_panel.line.width = Pt(1)
        
        right_box = slide.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.4))
        tf_right = right_box.text_frame
        tf_right.word_wrap = True
        add_bullets(tf_right, points_right)
    else:
        # Full Width Box
        full_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
        tf_full = full_box.text_frame
        tf_full.word_wrap = True
        add_bullets(tf_full, points_left)

def add_bullets(text_frame, bullet_list):
    first = True
    for header, desc in bullet_list:
        p = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False
        
        p.text = f"•  {header}"
        p.font.name = 'Outfit'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p.space_after = Pt(4)
        
        p_desc = text_frame.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = 'Plus Jakarta Sans'
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = COLOR_TEXT_MAIN
        p_desc.space_after = Pt(20)
        p_desc.level = 0
        # Indent description slightly
        p_desc.space_before = Pt(2)

# Slide 2: The Challenge & Vision
points_2_left = [
    ("Hackathon Goal", "Create a highly interactive, responsive FIFA World Cup 2026 dashboard using exclusively native front-end core technologies (HTML5, CSS3, Vanilla JS)."),
    ("Avoid Common Pitfalls", "Move away from standard 'student project' looks by utilizing high-end UI design systems: custom variables, detailed spacing, and rich typography."),
    ("Audience-First UX", "Provide clear readability, intuitive layouts, and GPU-accelerated micro-animations suitable for press, analysts, and fans alike.")
]
points_2_right = [
    ("Glassmorphism Aesthetic", "Frosted glass card layouts constructed using backdrop-filter blur parameters and subtle semi-transparent borders."),
    ("Rich Visual Branding", "Deep Navy gradient backdrop, subtle football pitch grids, gold rank badges, and glowing animations that mirror an official FIFA platform."),
    ("Dynamic Feel", "A responsive workspace that adjusts layout dimensions immediately to ensure content remains visible on all screen sizes.")
]

# Slide 3: Visual & UX Architecture
points_3_left = [
    ("Visual Theme Core", "Deep navy background coordinates paired with radial cyan/indigo glows and floating HTML particles that float upward automatically."),
    ("Frosted Glass Layouts", "Backdrop-filter: blur(12px) paired with custom box-shadow definitions for a premium layered 3D aesthetic."),
    ("First-Place Winner Accent", "Gold-tinted glass background, glowing golden border animation, crown badge, and hover elevate scale to draw eyes immediately.")
]
points_3_right = [
    ("CSS Custom Properties", "Centralized HSL color variables supporting complete Dark and Light theme modifications smoothly in real-time."),
    ("Interactive Flags CDN", "High-fidelity flags loaded dynamically from FlagCDN database ensuring lightweight performance and zero broken image paths."),
    ("Mobile Adaptation", "Responsive layout adjustments down to 320px mobile portrait. Stackable group cards, horizontally scrollable standings grids.")
]

# Slide 4: Key Functional Modules
points_4_left = [
    ("Sticky Navigation System", "Scroll-margin layouts supporting click-to-scroll navigation anchor targets smoothly, alongside immediate control actions."),
    ("Group Stage Explorer", "Accordion components allowing Group A, B, C, and D previews. Custom toggle engine ensures only one group card expands at a time."),
    ("Standings Analyzer Table", "Fully responsive data rows containing team flags, rankings, computed wins, losses, goals, and points statistics.")
]
points_4_right = [
    ("Fuzzy Search Filter", "Real-time keyup query checking that instantly hides non-matching table rows and Group Explorer team lists without page refresh."),
    ("Multi-Column Sorting", "Interactive column headers supporting descending/ascending toggles for Points, Wins, Played, Goals, and Goal Difference."),
    ("Recent Form Indicators", "Color-coded circular indicator badges (Green: Win, Yellow: Draw, Red: Loss) representing the last 3 match outcomes.")
]

# Slide 5: Technical Implementation
points_5_left = [
    ("Vanilla Core Stack", "Built without external CSS or JS frameworks (Bootstrap, Tailwind, React, etc.) ensuring lightning-fast client-side load speeds."),
    ("GPU-Accelerated Transition", "Transforms, scales, opacity and filter changes handled on the GPU for smooth 60fps renders on mobile and desktops."),
    ("Accessible ARIA Bindings", "Semantic outline structures, high-contrast text accessibility ratios, large touch target bounds, and visible focus outline states.")
]
points_5_right = [
    ("Dynamic Match Simulator", "Clicking the 'Refresh' button updates goals, recalculates points/goal differences, shifts rankings, and updates average stats dynamically."),
    ("Live Count Loader", "Custom requestAnimationFrame increments that animate statistics values (Teams, Matches) from zero on initial load."),
    ("Toast Notification System", "Glassmorphic alert prompts that slide up from the bottom corner to confirm actions and system updates.")
]

# Slide 6: Summary & Hackathon Credentials
points_6_full = [
    ("Clean Project Directory Structure", "Standard modular distribution containing: index.html (semantic outline), style.css (visual styles), script.js (logic engine), logo.png (branding), and README.md (instructions)."),
    ("Developer Credentials", "Built by Team 'FIFO' representing Country 'Uruguay'. Repository hosted officially at: https://github.com/notjamestbh"),
    ("Ready for Production", "Production-grade, optimized code suitable for winning hackathon selection rounds, fully tested with zero console logs or errors.")
]

# Build Presentation
create_title_slide()
create_content_slide("The Challenge & Vision", "Overview", points_2_left, points_2_right)
create_content_slide("Visual & UX Architecture", "Aesthetics", points_3_left, points_3_right)
create_content_slide("Key Functional Modules", "Interactive Features", points_4_left, points_4_right)
create_content_slide("Technical Implementation", "Performance & Code", points_5_left, points_5_right)
create_content_slide("Summary & Credentials", "Project Delivery", points_6_full)

# Save
prs.save(pptx_path)
print(f"Presentation successfully created at: {pptx_path}")
